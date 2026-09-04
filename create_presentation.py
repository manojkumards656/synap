import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_synap_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Design System Colors (Clean, Modern Light Theme)
    BG_COLOR = RGBColor(248, 250, 252)       # Soft off-white Slate 50
    CARD_BG = RGBColor(255, 255, 255)        # Pure white
    CARD_BORDER = RGBColor(226, 232, 240)    # Slate 200
    TEXT_MAIN = RGBColor(15, 23, 42)         # Slate 900
    TEXT_MUTED = RGBColor(100, 116, 139)     # Slate 500
    ACCENT_BLUE = RGBColor(37, 99, 235)      # Royal Blue 600
    ACCENT_CYAN = RGBColor(2, 132, 199)      # Sky 600
    ACCENT_GREEN = RGBColor(22, 163, 74)     # Emerald 600
    ACCENT_RED = RGBColor(220, 38, 38)       # Crimson 600
    CHIP_BG = RGBColor(239, 246, 255)        # Light Blue 50

    def apply_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return bg

    def add_header(slide, category, title, subtitle=None):
        cat_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.4))
        tf = cat_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = category.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.font.name = "Arial"

        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.95), Inches(11.333), Inches(0.8))
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_MAIN
        p2.font.name = "Arial"

        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.75), Inches(11.333), Inches(0.45))
            tf3 = sub_box.text_frame
            tf3.word_wrap = True
            tf3.margin_left = tf3.margin_top = tf3.margin_right = tf3.margin_bottom = 0
            p3 = tf3.paragraphs[0]
            p3.text = subtitle
            p3.font.size = Pt(14)
            p3.font.color.rgb = TEXT_MUTED
            p3.font.name = "Arial"

    def add_card(slide, left, top, width, height, border_color=CARD_BORDER, bg_color=CARD_BG):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        return card

    # =========================================================================
    # SLIDE 1: Title Slide (Simple, Clear & Impactful)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_background(slide1)

    pill = add_card(slide1, Inches(1.0), Inches(1.8), Inches(3.6), Inches(0.45), border_color=ACCENT_BLUE, bg_color=CHIP_BG)
    tf_pill = pill.text_frame
    tf_pill.word_wrap = False
    p_pill = tf_pill.paragraphs[0]
    p_pill.text = "🛡️ REAL-TIME SCAM PROTECTION FOR EVERYONE"
    p_pill.font.size = Pt(10)
    p_pill.font.bold = True
    p_pill.font.color.rgb = ACCENT_BLUE
    p_pill.alignment = PP_ALIGN.CENTER

    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(1.4))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SYNAP"
    p.font.size = Pt(76)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    p.font.name = "Arial"

    sub_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(10.5), Inches(1.2))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Protecting Bank Accounts from Phone Call Scams in Real Time"
    p_sub.font.size = Pt(24)
    p_sub.font.bold = True
    p_sub.font.color.rgb = ACCENT_BLUE
    p_sub.font.name = "Arial"

    desc_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(10.0), Inches(1.0))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "When fraudsters trick victims over a phone call into sending money, passwords and OTPs cannot help.\nSynap detects the natural signs of panic in how a user touches their screen — and pauses the transfer."
    p_desc.font.size = Pt(15)
    p_desc.font.color.rgb = TEXT_MUTED
    p_desc.font.name = "Arial"

    # =========================================================================
    # SLIDE 2: The Real-World Scam (Easy to Understand for Everyone)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_background(slide2)
    add_header(slide2, "The Real-World Threat", "How Scammers Trick People Over the Phone", "Criminals no longer hack devices. They manipulate real people in real time.")

    card_width = Inches(3.55)
    card_height = Inches(4.2)
    gap = Inches(0.34)
    start_x = Inches(1.0)
    top_y = Inches(2.4)

    # Card 1: The Phone Call
    add_card(slide2, start_x, top_y, card_width, card_height)
    b1 = slide2.shapes.add_textbox(start_x + Inches(0.3), top_y + Inches(0.4), card_width - Inches(0.6), card_height - Inches(0.8))
    t1 = b1.text_frame
    t1.word_wrap = True
    p1 = t1.paragraphs[0]
    p1.text = "01"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_BLUE

    p1_sub = t1.add_paragraph()
    p1_sub.text = "The Panic Phone Call"
    p1_sub.font.size = Pt(18)
    p1_sub.font.bold = True
    p1_sub.font.color.rgb = TEXT_MAIN

    p1_body = t1.add_paragraph()
    p1_body.text = "\n• Scammers call posing as police, tax, or bank fraud officials.\n• They create severe fear & urgency: 'Your account is compromised!'\n• They order the victim to move money to a 'safe reserve account'."
    p1_body.font.size = Pt(13)
    p1_body.font.color.rgb = TEXT_MUTED

    # Card 2: Traditional Security Fails
    x2 = start_x + card_width + gap
    add_card(slide2, x2, top_y, card_width, card_height)
    b2 = slide2.shapes.add_textbox(x2 + Inches(0.3), top_y + Inches(0.4), card_width - Inches(0.6), card_height - Inches(0.8))
    t2 = b2.text_frame
    t2.word_wrap = True
    p2 = t2.paragraphs[0]
    p2.text = "02"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_RED

    p2_sub = t2.add_paragraph()
    p2_sub.text = "Why Passwords Fail"
    p2_sub.font.size = Pt(18)
    p2_sub.font.bold = True
    p2_sub.font.color.rgb = TEXT_MAIN

    p2_body = t2.add_paragraph()
    p2_body.text = "\n• The victim uses their own phone and enters their own PIN.\n• They pass FaceID / fingerprint voluntarily.\n• The bank's security sees normal logins and lets the money go.\n• Standard fraud filters are completely blind."
    p2_body.font.size = Pt(13)
    p2_body.font.color.rgb = TEXT_MUTED

    # Card 3: The Human Impact
    x3 = x2 + card_width + gap
    add_card(slide2, x3, top_y, card_width, card_height)
    b3 = slide2.shapes.add_textbox(x3 + Inches(0.3), top_y + Inches(0.4), card_width - Inches(0.6), card_height - Inches(0.8))
    t3 = b3.text_frame
    t3.word_wrap = True
    p3 = t3.paragraphs[0]
    p3.text = "03"
    p3.font.size = Pt(28)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_GREEN

    p3_sub = t3.add_paragraph()
    p3_sub.text = "The Devastating Loss"
    p3_sub.font.size = Pt(18)
    p3_sub.font.bold = True
    p3_sub.font.color.rgb = TEXT_MAIN

    p3_body = t3.add_paragraph()
    p3_body.text = "\n• Senior citizens and first-time digital users lose life savings in minutes.\n• Victims feel helpless, embarrassed, and financially ruined.\n• Over ₹4,000+ Crores stolen every year via these phone scams."
    p3_body.font.size = Pt(13)
    p3_body.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 3: The Big Idea (Simple Contrast)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_background(slide3)
    add_header(slide3, "The Big Idea", "Don't Check WHO is Typing. Check HOW They Are Typing.", "A scammer can force you to type your password, but they cannot fake calm behavior.")

    panel_w = Inches(5.45)
    panel_h = Inches(4.3)

    # Left: What Banks Do Today
    add_card(slide3, Inches(1.0), Inches(2.3), panel_w, panel_h)
    p_left = slide3.shapes.add_textbox(Inches(1.4), Inches(2.6), panel_w - Inches(0.8), panel_h - Inches(0.6))
    tf_l = p_left.text_frame
    tf_l.word_wrap = True
    h1 = tf_l.paragraphs[0]
    h1.text = "WHAT CURRENT APPS DO"
    h1.font.size = Pt(13)
    h1.font.bold = True
    h1.font.color.rgb = ACCENT_RED

    h1_title = tf_l.add_paragraph()
    h1_title.text = "Checking Identity (Blind to Fear)"
    h1_title.font.size = Pt(21)
    h1_title.font.bold = True
    h1_title.font.color.rgb = TEXT_MAIN

    h1_body = tf_l.add_paragraph()
    h1_body.text = "\n• Asks: 'Do you know the password?' (Yes)\n• Asks: 'Is this your fingerprint?' (Yes)\n• Asks: 'Is this your usual phone?' (Yes)\n\n❌ Flaw: All of these pass even when the victim is being threatened on a live phone call."
    h1_body.font.size = Pt(13)
    h1_body.font.color.rgb = TEXT_MUTED

    # Right: What Synap Does
    add_card(slide3, Inches(6.88), Inches(2.3), panel_w, panel_h, border_color=ACCENT_BLUE)
    p_right = slide3.shapes.add_textbox(Inches(7.28), Inches(2.6), panel_w - Inches(0.8), panel_h - Inches(0.6))
    tf_r = p_right.text_frame
    tf_r.word_wrap = True
    h2 = tf_r.paragraphs[0]
    h2.text = "WHAT SYNAP DOES"
    h2.font.size = Pt(13)
    h2.font.bold = True
    h2.font.color.rgb = ACCENT_BLUE

    h2_title = tf_r.add_paragraph()
    h2_title.text = "Checking Behavior (Catches Panic)"
    h2_title.font.size = Pt(21)
    h2_title.font.bold = True
    h2_title.font.color.rgb = TEXT_MAIN

    h2_body = tf_r.add_paragraph()
    h2_body.text = "\n• Detects long pauses while listening to dictation\n• Detects panicked, rapid typing bursts\n• Detects shaky, uneven finger pressure\n• Checks if a live phone call is running right now\n\n✅ Result: Stops the money before it leaves the account."
    h2_body.font.size = Pt(13)
    h2_body.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 4: How It Works (Simple Signals, Zero Jargon)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_background(slide4)
    add_header(slide4, "The Technology Made Simple", "How Synap Senses Trouble (Without Reading Your Data)", "Synap looks at the natural physics of touch — completely private and safe.")

    gw = Inches(3.55)
    gh = Inches(1.85)
    gx_gap = Inches(0.34)
    gy_gap = Inches(0.25)

    signals = [
        ("01", "Listening Pauses", "Dictation Check", "Victim stops typing for 2–3 seconds while listening to the scammer give the next number."),
        ("02", "Panic Bursts", "Rush Typing", "Victim rushes to enter digits quickly under pressure, hitting keys in rapid bursts."),
        ("03", "Shaky Press Durations", "Tremor Detection", "Fingers linger irregularly on keys due to nervous tension and anxiety."),
        ("04", "Uneven Rhythm", "Broken Flow", "Normal typing is steady. Coerced typing starts and stops unpredictably."),
        ("05", "Chaotic Finger Path", "Wavering Motion", "Fingers waver across the screen instead of moving in clean, direct paths."),
        ("06", "Live Phone Call Check", "The Key Trigger", "Verifies if the person is on an active phone call during the high-value transfer."),
    ]

    for idx, (num, title, subtitle, desc) in enumerate(signals):
        col = idx % 3
        row = idx // 3
        cx = Inches(1.0) + col * (gw + gx_gap)
        cy = Inches(2.2) + row * (gh + gy_gap)

        add_card(slide4, cx, cy, gw, gh)
        s_box = slide4.shapes.add_textbox(cx + Inches(0.25), cy + Inches(0.18), gw - Inches(0.5), gh - Inches(0.36))
        stf = s_box.text_frame
        stf.word_wrap = True

        sp1 = stf.paragraphs[0]
        sp1.text = f"{num} • {title}"
        sp1.font.size = Pt(14)
        sp1.font.bold = True
        sp1.font.color.rgb = TEXT_MAIN

        sp2 = stf.add_paragraph()
        sp2.text = f"{subtitle}: {desc}"
        sp2.font.size = Pt(11)
        sp2.font.color.rgb = TEXT_MUTED

    b_pill = add_card(slide4, Inches(1.0), Inches(6.35), Inches(11.333), Inches(0.55), border_color=ACCENT_GREEN, bg_color=CARD_BG)
    btf = b_pill.text_frame
    btf.word_wrap = False
    bp = btf.paragraphs[0]
    bp.text = "🔒 100% PRIVATE: Synap never reads messages, contacts, passwords, or calls. It only measures the timing of taps."
    bp.font.size = Pt(11)
    bp.font.bold = True
    bp.font.color.rgb = ACCENT_GREEN
    bp.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 5: How It Directly Helps Users (Elderly Friendly Focus)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_background(slide5)
    add_header(slide5, "Human-First Innovation", "How Synap Protects Everyday Users & Seniors", "Built like Google Pay — simple, reassuring, and stress-free.")

    p5_w = Inches(3.55)
    p5_h = Inches(4.3)

    # Benefit 1: Protection Under Pressure
    add_card(slide5, Inches(1.0), Inches(2.3), p5_w, p5_h)
    b_u1 = slide5.shapes.add_textbox(Inches(1.25), Inches(2.6), p5_w - Inches(0.5), p5_h - Inches(0.6))
    tu1 = b_u1.text_frame
    tu1.word_wrap = True
    up1 = tu1.paragraphs[0]
    up1.text = "TIME TO BREATHE"
    up1.font.size = Pt(12)
    up1.font.bold = True
    up1.font.color.rgb = ACCENT_BLUE

    up1_t = tu1.add_paragraph()
    up1_t.text = "The 15-Minute Safety Hold"
    up1_t.font.size = Pt(18)
    up1_t.font.bold = True
    up1_t.font.color.rgb = TEXT_MAIN

    up1_b = tu1.add_paragraph()
    up1_b.text = "\n• Scammers rely on rush and fear.\n• Synap introduces a mandatory 15-minute cooling-off window.\n• Breaks the psychological hold of the caller.\n• Gives the victim time to call family or verify with the bank."
    up1_b.font.size = Pt(12)
    up1_b.font.color.rgb = TEXT_MUTED

    # Benefit 2: Senior-Friendly Experience
    add_card(slide5, Inches(4.89), Inches(2.3), p5_w, p5_h)
    b_u2 = slide5.shapes.add_textbox(Inches(5.14), Inches(2.6), p5_w - Inches(0.5), p5_h - Inches(0.6))
    tu2 = b_u2.text_frame
    tu2.word_wrap = True
    up2 = tu2.paragraphs[0]
    up2.text = "ELDERLY-FRIENDLY DESIGN"
    up2.font.size = Pt(12)
    up2.font.bold = True
    up2.font.color.rgb = ACCENT_GREEN

    up2_t = tu2.add_paragraph()
    up2_t.text = "Clarity, Not Confusion"
    up2_t.font.size = Pt(18)
    up2_t.font.bold = True
    up2_t.font.color.rgb = TEXT_MAIN

    up2_b = tu2.add_paragraph()
    up2_b.text = "\n• Large Keypad (64dp buttons) for easy tapping with tremors.\n• Spoken Words: '50000' shows as 'Fifty Thousand Rupees' so seniors never count zeros.\n• Reassuring Banner: 'Do not worry, your money is safe in your account'."
    up2_b.font.size = Pt(12)
    up2_b.font.color.rgb = TEXT_MUTED

    # Benefit 3: Zero Complications
    add_card(slide5, Inches(8.78), Inches(2.3), p5_w, p5_h)
    b_u3 = slide5.shapes.add_textbox(Inches(9.03), Inches(2.6), p5_w - Inches(0.5), p5_h - Inches(0.6))
    tu3 = b_u3.text_frame
    tu3.word_wrap = True
    up3 = tu3.paragraphs[0]
    up3.text = "STRESS-FREE RECOVERY"
    up3.font.size = Pt(12)
    up3.font.bold = True
    up3.font.color.rgb = ACCENT_RED

    up3_t = tu3.add_paragraph()
    up3_t.text = "1-Tap Cancel & Help"
    up3_t.font.size = Pt(18)
    up3_t.font.bold = True
    up3_t.font.color.rgb = TEXT_MAIN

    up3_b = tu3.add_paragraph()
    up3_b.text = "\n• Prominent Green Button: 'Cancel Transfer & Keep My Money Safe'.\n• Direct 24/7 Senior Citizen Fraud Helpline shortcut.\n• No annoying lockouts or forgotten password resets for everyday users."
    up3_b.font.size = Pt(12)
    up3_b.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 6: Live Demo Comparison (Simple A/B Test)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_background(slide6)
    add_header(slide6, "The Live Demonstration", "Two Everyday Situations. One Smart App.", "How Synap tells the difference between buying groceries and being scammed.")

    f_w = Inches(5.45)
    f_h = Inches(4.5)

    # Scenario A: Calm
    add_card(slide6, Inches(1.0), Inches(2.3), f_w, f_h, border_color=ACCENT_GREEN)
    box_a = slide6.shapes.add_textbox(Inches(1.35), Inches(2.6), f_w - Inches(0.7), f_h - Inches(0.6))
    tfa = box_a.text_frame
    tfa.word_wrap = True

    ap1 = tfa.paragraphs[0]
    ap1.text = "SCENARIO 1: NORMAL DAILY TRANSFER"
    ap1.font.size = Pt(13)
    ap1.font.bold = True
    ap1.font.color.rgb = ACCENT_GREEN

    ap2 = tfa.add_paragraph()
    ap2.text = "Instant, Frictionless Payment"
    ap2.font.size = Pt(21)
    ap2.font.bold = True
    ap2.font.color.rgb = TEXT_MAIN

    ap3 = tfa.add_paragraph()
    ap3.text = "\n• Phone is NOT on an active call.\n• User types calmly and at their own relaxed pace.\n• Duress score stays safely GREEN (under 0.30).\n• Payment clears in under 1 second without any delay.\n• Perfect for everyday shopping and family transfers."
    ap3.font.size = Pt(13)
    ap3.font.color.rgb = TEXT_MUTED

    # Scenario B: Coerced
    add_card(slide6, Inches(6.88), Inches(2.3), f_w, f_h, border_color=ACCENT_RED)
    box_b = slide6.shapes.add_textbox(Inches(7.23), Inches(2.6), f_w - Inches(0.7), f_h - Inches(0.6))
    tfb = box_b.text_frame
    tfb.word_wrap = True

    bp1 = tfb.paragraphs[0]
    bp1.text = "SCENARIO 2: TRICKED OVER A PHONE CALL"
    bp1.font.size = Pt(13)
    bp1.font.bold = True
    bp1.font.color.rgb = ACCENT_RED

    bp2 = tfb.add_paragraph()
    bp2.text = "Protected with a 15-Minute Pause"
    bp2.font.size = Pt(21)
    bp2.font.bold = True
    bp2.font.color.rgb = TEXT_MAIN

    bp3 = tfb.add_paragraph()
    bp3.text = "\n• A phone call is ACTIVE during the transfer.\n• User pauses for 3 seconds listening to instructions, then rushes.\n• Shaky finger pressure pushes duress score into the RED (0.85).\n• App instantly pauses the transfer for 15 minutes.\n• The money never leaves the account. Life savings saved!"
    bp3.font.size = Pt(13)
    bp3.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 7: Why Synap Wins (Summary & Impact)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_background(slide7)
    add_header(slide7, "The Big Picture", "A Safer Digital Banking Experience for India", "Protecting families and banks without compromising speed or privacy.")

    s_w = Inches(3.55)
    s_h = Inches(4.3)

    vals = [
        ("Peace of Mind for Families", "Protecting Seniors", "Seniors can use digital banking with confidence, knowing an invisible safety net stops coercion before money is lost."),
        ("No Extra Friction", "Lightning-Fast (<0.5ms)", "Legitimate UPI payments are never delayed. No annoying puzzles, CAPTCHAs, or extra codes to memorize."),
        ("Bank Loss Prevention", "Direct Liability Shield", "Banks save crores in scam reimbursements and regulatory penalties under new digital payment safety guidelines.")
    ]

    for idx, (title, highlight, body) in enumerate(vals):
        vx = Inches(1.0) + idx * (s_w + Inches(0.34))
        vy = Inches(2.3)
        add_card(slide7, vx, vy, s_w, s_h)

        v_box = slide7.shapes.add_textbox(vx + Inches(0.3), vy + Inches(0.4), s_w - Inches(0.6), s_h - Inches(0.8))
        vtf = v_box.text_frame
        vtf.word_wrap = True

        vp1 = vtf.paragraphs[0]
        vp1.text = highlight.upper()
        vp1.font.size = Pt(12)
        vp1.font.bold = True
        vp1.font.color.rgb = ACCENT_BLUE

        vp2 = vtf.add_paragraph()
        vp2.text = title
        vp2.font.size = Pt(19)
        vp2.font.bold = True
        vp2.font.color.rgb = TEXT_MAIN

        vp3 = vtf.add_paragraph()
        vp3.text = f"\n{body}"
        vp3.font.size = Pt(13)
        vp3.font.color.rgb = TEXT_MUTED

    output_path = os.path.abspath("synap_presentation.pptx")
    prs.save(output_path)
    print(f"Updated presentation saved successfully at: {output_path}")

if __name__ == "__main__":
    create_synap_presentation()
